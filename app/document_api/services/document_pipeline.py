from __future__ import annotations

import hashlib
import json
import math
import mimetypes
import time
import re
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Iterable

from psycopg2 import sql as psql_sql

from app.document_api.core.config import default_mcp_km_number_of_results
from app.document_api.services.office_legacy_convert import LEGACY_OFFICE_EXTENSIONS


def _safe_storage_name(name: str) -> str:
    base = Path((name or "").strip().replace("\\", "/")).name
    stem = Path(base).stem
    ext = Path(base).suffix.lower()

    safe_stem = re.sub(r"[^A-Za-z0-9._-]+", "-", stem).strip("._-").lower() or "file"
    safe_ext = ext if re.fullmatch(r"\.[a-z0-9]{1,10}", ext or "") else ""
    digest = hashlib.md5(base.encode("utf-8")).hexdigest()[:10]
    return f"{safe_stem}-{digest}{safe_ext}"


def _vector_table_ident(table_name: str | None) -> str:
    """ชื่อตาราง chunk/vector จาก env — ต้องเป็น PostgreSQL identifier ปกติ (ไม่มี quote/ช่องว่าง)."""
    t = (table_name or "").strip()
    if not re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]{0,62}", t):
        raise ValueError(
            "SUPABASE_TABLE_NAME / sb.table_name ต้องเป็น PostgreSQL identifier "
            "(ขึ้นต้นด้วยตัวอักษรหรือ _ แล้วตามด้วยตัวอักษร ตัวเลข _ ความยาวไม่เกิน 63): "
            f"ได้รับ {table_name!r}"
        )
    return t


def _vector_rows_for_source_sql(doc_name: str | None) -> tuple[str, list[str]]:
    """Scope vector rows to one document set when document_name is known."""
    doc = (doc_name or "").strip()
    if not doc:
        return "metadata->>'source_filename' = %s", []
    return (
        "metadata->>'source_filename' = %s"
        " AND (document_name = %s OR (metadata->>'document_name') = %s)",
        [doc, doc],
    )


# ===========================================================================
# ── BASE CONVERTER
# ===========================================================================

@dataclass
class ImageAsset:
    name: str
    path: Path
    rel_path: str
    alt_text: str | None = None
    caption: str | None = None
    page: int | None = None


@dataclass
class ConversionResult:
    markdown: str
    images: list[ImageAsset] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    source_filename: str = ""


class BaseConverter:
    SUPPORTED_EXTENSIONS: tuple[str, ...] = ()

    def __init__(self) -> None:
        self._image_counter = 0

    def supports(self, suffix: str) -> bool:
        return suffix.lower() in self.SUPPORTED_EXTENSIONS

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        raise NotImplementedError

    @staticmethod
    def _ensure_image_dir(output_dir: Path) -> Path:
        img_dir = output_dir / "images"
        img_dir.mkdir(parents=True, exist_ok=True)
        return img_dir

    def _save_image_bytes(
        self,
        data: bytes,
        output_dir: Path,
        ext: str = "png",
        page: int | None = None,
    ) -> ImageAsset:
        img_dir = self._ensure_image_dir(output_dir)
        self._image_counter += 1
        digest = hashlib.md5(data).hexdigest()[:8]
        ext = ext.lstrip(".").lower() or "png"
        name = f"img_{self._image_counter:03d}_{digest}.{ext}"
        path = img_dir / name
        path.write_bytes(data)
        return ImageAsset(name=name, path=path, rel_path=f"images/{name}", page=page)

    @staticmethod
    def _md_image(asset: ImageAsset) -> str:
        return f"![{asset.alt_text or asset.name}]({asset.rel_path})"

    @staticmethod
    def join_blocks(blocks: Iterable[str]) -> str:
        return "\n\n".join(b.strip() for b in blocks if b and b.strip())


# ===========================================================================
# ── PDF
# ===========================================================================


class PDFConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".pdf",)

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        import fitz

        doc = fitz.open(file_path)
        try:
            blocks: list[str] = [f"# {file_path.stem}"]
            images: list[ImageAsset] = []
            metadata = {
                "page_count": doc.page_count,
                "title": doc.metadata.get("title") if doc.metadata else None,
                "author": doc.metadata.get("author") if doc.metadata else None,
            }
            for i in range(doc.page_count):
                page = doc.load_page(i)
                pno = i + 1
                blocks.append(f"## หน้า {pno}")
                text = page.get_text("text").strip()
                if text:
                    blocks.append(text)
                for info in page.get_images(full=True):
                    try:
                        base = doc.extract_image(info[0])
                    except Exception:
                        continue
                    img_bytes = base.get("image")
                    if not img_bytes:
                        continue
                    asset = self._save_image_bytes(
                        img_bytes,
                        output_dir,
                        ext=base.get("ext", "png"),
                        page=pno,
                    )
                    images.append(asset)
                    blocks.append(self._md_image(asset))
            return ConversionResult(
                markdown=self.join_blocks(blocks),
                images=images,
                metadata=metadata,
                source_filename=file_path.name,
            )
        finally:
            doc.close()


# ===========================================================================
# ── DOCX
# ===========================================================================


class DOCXConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".docx", ".docm")

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        from docx import Document
        from docx.document import Document as _Doc
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        doc: _Doc = Document(str(file_path))
        blocks: list[str] = [f"# {file_path.stem}"]
        images: list[ImageAsset] = []
        for child in doc.element.body.iterchildren():
            tag = child.tag.split("}")[-1]
            if tag == "p":
                md = self._para_to_md(Paragraph(child, doc), doc, output_dir, images)
                if md:
                    blocks.append(md)
            elif tag == "tbl":
                blocks.append(self._table_to_md(Table(child, doc)))
        return ConversionResult(
            markdown=self.join_blocks(blocks),
            images=images,
            metadata={"paragraph_count": len(doc.paragraphs)},
            source_filename=file_path.name,
        )

    def _para_to_md(self, para, doc, output_dir, images) -> str:
        text = para.text.strip()
        embedded = self._para_images(para, doc, output_dir, images)
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            try:
                level = max(1, min(int(style.replace("heading", "").strip() or "2"), 6))
            except ValueError:
                level = 2
            return self.join_blocks([f"{'#' * level} {text}" if text else "", *embedded])
        if style in {"list paragraph", "list bullet"}:
            return self.join_blocks([f"- {text}" if text else "", *embedded])
        if not text and not embedded:
            return ""
        return self.join_blocks([text, *embedded])

    def _para_images(self, para, doc, output_dir, images) -> list[str]:
        from docx.oxml.ns import qn

        out: list[str] = []
        for blip in para._p.iter(qn("a:blip")):
            rid = blip.get(qn("r:embed"))
            if not rid:
                continue
            try:
                part = doc.part.related_parts[rid]
            except KeyError:
                continue
            ext = part.partname.ext.lstrip(".") if hasattr(part.partname, "ext") else "png"
            asset = self._save_image_bytes(part.blob, output_dir, ext=ext)
            images.append(asset)
            out.append(self._md_image(asset))
        return out

    @staticmethod
    def _table_to_md(table) -> str:
        rows = [[c.text.strip().replace("\n", " ") for c in r.cells] for r in table.rows]
        if not rows:
            return ""
        w = len(rows[0])

        def line(cells):
            p = cells[:w] + [""] * max(0, w - len(cells))
            return "| " + " | ".join(p) + " |"

        return "\n".join([line(rows[0]), line(["---"] * w)] + [line(r) for r in rows[1:]])


# ===========================================================================
# ── PPTX
# ===========================================================================


class PPTXConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".pptx", ".pptm")

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(file_path))
        blocks: list[str] = [f"# {file_path.stem}"]
        images: list[ImageAsset] = []
        for idx, slide in enumerate(prs.slides, 1):
            sb: list[str] = [f"## สไลด์ {idx}"]
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                t = slide.shapes.title.text_frame.text.strip()
                if t:
                    sb.append(f"### {t}")
            for shape in slide.shapes:
                self._shape(shape, output_dir, sb, images, idx, MSO_SHAPE_TYPE)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    sb.append("> **Notes:** " + notes.replace("\n", " "))
            blocks.append(self.join_blocks(sb))
        return ConversionResult(
            markdown=self.join_blocks(blocks),
            images=images,
            metadata={"slide_count": len(prs.slides)},
            source_filename=file_path.name,
        )

    def _shape(self, shape, output_dir, sb, images, page, MSO):
        if shape.shape_type == MSO.GROUP:
            for s in shape.shapes:
                self._shape(s, output_dir, sb, images, page, MSO)
        elif shape.shape_type == MSO.PICTURE:
            try:
                ext = (shape.image.ext or "png").lstrip(".")
                asset = self._save_image_bytes(shape.image.blob, output_dir, ext=ext, page=page)
                images.append(asset)
                sb.append(self._md_image(asset))
            except Exception:
                pass
        elif shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                sb.append(t)
        elif shape.has_table:
            rows = [[c.text.strip().replace("\n", " ") for c in r.cells] for r in shape.table.rows]
            if rows:
                w = len(rows[0])

                def line(cells):
                    p = cells[:w] + [""] * max(0, w - len(cells))
                    return "| " + " | ".join(p) + " |"

                sb.append("\n".join([line(rows[0]), line(["---"] * w)] + [line(r) for r in rows[1:]]))


# ===========================================================================
# ── XLSX helpers
# ===========================================================================


def _detect_ext(data: bytes) -> str:
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return "png"
    if data[:3] == b"\xff\xd8\xff":
        return "jpg"
    if data[:6] in (b"GIF87a", b"GIF89a"):
        return "gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "webp"
    if data[:2] == b"BM":
        return "bmp"
    if data[:4] in (b"II*\x00", b"MM\x00*"):
        return "tiff"
    return "png"


def _img_bytes(img_obj) -> bytes | None:
    fn = getattr(img_obj, "_data", None)
    if callable(fn):
        try:
            d = fn()
            if d:
                return d
        except Exception:
            pass
    ref = getattr(img_obj, "ref", None)
    if ref and hasattr(ref, "read"):
        try:
            ref.seek(0)
            d = ref.read()
            if d:
                return d
        except Exception:
            pass
    if ref and hasattr(ref, "getvalue"):
        try:
            d = ref.getvalue()
            if d:
                return d
        except Exception:
            pass
    pil = getattr(img_obj, "image", None)
    if pil:
        try:
            import io

            buf = io.BytesIO()
            pil.save(buf, format=getattr(pil, "format", None) or "PNG")
            d = buf.getvalue()
            if d:
                return d
        except Exception:
            pass
    return None


# ===========================================================================
# ── XLSX
# ===========================================================================


class XLSXConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".xlsx", ".xlsm")

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        import warnings

        import openpyxl
        from openpyxl.utils import get_column_letter

        blocks: list[str] = [f"# {file_path.stem}"]
        images: list[ImageAsset] = []
        skipped = 0
        with warnings.catch_warnings():
            warnings.filterwarnings(
                "ignore",
                message="DrawingML support is incomplete",
                category=UserWarning,
                module="openpyxl.reader.drawings",
            )
            wb = openpyxl.load_workbook(file_path, data_only=True)

        for sn in wb.sheetnames:
            ws = wb[sn]
            sb: list[str] = [f"## Sheet: {sn}"]
            tabular = sum(
                1
                for row in ws.iter_rows(min_row=1, max_row=min(20, ws.max_row), values_only=True)
                if sum(1 for v in row if v is not None and str(v).strip()) >= 3
            )
            sb.append(self._to_table(ws) if tabular >= 3 else self._to_text(ws))

            anchors = []
            for img_obj in getattr(ws, "_images", []):
                data = _img_bytes(img_obj)
                if not data:
                    skipped += 1
                    continue
                anchor = getattr(img_obj, "anchor", None)
                fr = fc = tr = tc = 0
                if anchor:
                    try:
                        _f = getattr(anchor, "_from", None)
                        if _f:
                            fr, fc = int(_f.row) + 1, int(_f.col) + 1
                        _t = getattr(anchor, "to", None)
                        if _t:
                            tr, tc = int(_t.row) + 1, int(_t.col) + 1
                        else:
                            tr, tc = fr, fc
                    except Exception:
                        pass
                ext = _detect_ext(data)
                asset = self._save_image_bytes(data, output_dir, ext=ext, page=fr or None)
                cell_ref = (
                    f"{get_column_letter(fc)}{fr}:{get_column_letter(max(tc, fc))}{max(tr, fr)}"
                    if fr and fc
                    else "(no anchor)"
                )
                asset.alt_text = f"{sn} {cell_ref}"
                asset.caption = asset.alt_text
                anchors.append((fr, fc, tr, tc, asset))

            anchors.sort(key=lambda t: (t[0], t[1]))
            if anchors:
                sb.append(f"### รูปภาพใน Sheet: {sn}")
                for fr, fc, tr, tc, asset in anchors:
                    images.append(asset)
                    sb.append(f"**{asset.alt_text}**\n\n{self._md_image(asset)}")

            blocks.append(self.join_blocks(sb))

        try:
            wb.close()
        except Exception:
            pass

        return ConversionResult(
            markdown=self.join_blocks(blocks),
            images=images,
            metadata={
                "sheet_count": len(wb.sheetnames),
                "image_count": len(images),
                "skipped_images": skipped,
            },
            source_filename=file_path.name,
        )

    @staticmethod
    def _rows_to_table(rows: list[list[str]]) -> str:
        if not rows:
            return ""
        w = len(rows[0])

        def line(cells: list[str]) -> str:
            padded = cells[:w] + [""] * max(0, w - len(cells))
            return "| " + " | ".join(c.replace("|", "\\|").replace("\n", " ") for c in padded) + " |"

        return "\n".join([line(rows[0]), line(["---"] * w)] + [line(r) for r in rows[1:]])

    @staticmethod
    def _rows_to_text(rows: list[list[str]]) -> str:
        if not rows:
            return ""
        lines: list[str] = []
        for vals in rows:
            cleaned = [v for v in vals if v]
            if not cleaned:
                continue
            lines.append(" | ".join(cleaned) if len(cleaned) > 1 else cleaned[0])
        return "\n\n".join(lines)

    @staticmethod
    def _to_table(ws) -> str:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(v).strip() if v is not None else "" for v in row]
            while cells and not cells[-1]:
                cells.pop()
            if cells:
                rows.append(cells)
        return XLSXConverter._rows_to_table(rows)

    @staticmethod
    def _to_text(ws) -> str:
        row_data: dict[int, list[str]] = {}
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None:
                    continue
                val = str(cell.value).strip()
                if val:
                    row_data.setdefault(cell.row, []).append(val)
        rows = [row_data[rn] for rn in sorted(row_data)]
        return XLSXConverter._rows_to_text(rows)


class XLSConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".xls",)

    @staticmethod
    def is_available() -> bool:
        try:
            import xlrd  # noqa: F401

            return True
        except ImportError:
            return False

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        import xlrd

        blocks: list[str] = [f"# {file_path.stem}"]
        book = xlrd.open_workbook(str(file_path))
        for sheet_name in book.sheet_names():
            sheet = book.sheet_by_name(sheet_name)
            sb: list[str] = [f"## Sheet: {sheet_name}"]
            rows: list[list[str]] = []
            for row_idx in range(sheet.nrows):
                full_row = [
                    str(sheet.cell_value(row_idx, col_idx)).strip()
                    if sheet.cell_value(row_idx, col_idx) is not None
                    else ""
                    for col_idx in range(sheet.ncols)
                ]
                while full_row and not full_row[-1]:
                    full_row.pop()
                if full_row:
                    rows.append(full_row)
            tabular = sum(1 for row in rows[:20] if sum(1 for v in row if v) >= 3)
            sb.append(
                XLSXConverter._rows_to_table(rows)
                if tabular >= 3
                else XLSXConverter._rows_to_text(rows)
            )
            blocks.append(self.join_blocks(sb))

        return ConversionResult(
            markdown=self.join_blocks(blocks),
            images=[],
            metadata={"sheet_count": len(book.sheet_names()), "image_count": 0, "format": "xls"},
            source_filename=file_path.name,
        )


# ===========================================================================
# ── CSV
# ===========================================================================


class CSVConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".csv",)

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        import pandas as pd

        try:
            df = pd.read_csv(file_path)
        except UnicodeDecodeError:
            df = pd.read_csv(file_path, encoding="latin-1")
        df = df.fillna("")
        return ConversionResult(
            markdown=self.join_blocks([f"# {file_path.stem}", df.to_markdown(index=False)]),
            images=[],
            metadata={"rows": int(df.shape[0]), "cols": int(df.shape[1])},
            source_filename=file_path.name,
        )


# ===========================================================================
# ── HTML / TEXT / IMAGE
# ===========================================================================


class HTMLConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".html", ".htm")

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        import base64
        from urllib.parse import urlparse

        import chardet
        from bs4 import BeautifulSoup
        from markdownify import markdownify as md_convert

        raw = file_path.read_bytes()
        enc = chardet.detect(raw).get("encoding") or "utf-8"
        soup = BeautifulSoup(raw.decode(enc, errors="ignore"), "lxml")
        images: list[ImageAsset] = []
        for img in soup.find_all("img"):
            src = img.get("src", "")
            if src.startswith("data:image/"):
                try:
                    header, b64 = src.split(",", 1)
                    ext = header.split("/")[1].split(";")[0]
                    asset = self._save_image_bytes(base64.b64decode(b64), output_dir, ext=ext)
                    images.append(asset)
                    img["src"] = asset.rel_path
                except Exception:
                    continue
            elif not urlparse(src).scheme:
                candidate = (file_path.parent / src).resolve()
                if candidate.exists():
                    ext = candidate.suffix.lstrip(".") or "png"
                    asset = self._save_image_bytes(candidate.read_bytes(), output_dir, ext=ext)
                    images.append(asset)
                    img["src"] = asset.rel_path
        return ConversionResult(
            markdown=md_convert(str(soup), heading_style="ATX").strip(),
            images=images,
            metadata={"title": soup.title.string if soup.title else None},
            source_filename=file_path.name,
        )


class TextConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".txt", ".md", ".markdown", ".log")

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        import chardet

        raw = file_path.read_bytes()
        enc = chardet.detect(raw).get("encoding") or "utf-8"
        text = raw.decode(enc, errors="ignore")
        md_text = text if file_path.suffix.lower() in (".md", ".markdown") else f"# {file_path.stem}\n\n{text}"
        return ConversionResult(
            markdown=md_text.strip(),
            images=[],
            metadata={"encoding": enc, "chars": len(text)},
            source_filename=file_path.name,
        )


class ImageFileConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff")

    def convert(self, file_path: Path, output_dir: Path) -> ConversionResult:
        ext = file_path.suffix.lstrip(".") or "png"
        asset = self._save_image_bytes(file_path.read_bytes(), output_dir, ext=ext)
        return ConversionResult(
            markdown=self.join_blocks([f"# {file_path.stem}", self._md_image(asset)]),
            images=[asset],
            metadata={},
            source_filename=file_path.name,
        )


# ===========================================================================
# ── CONVERTER REGISTRY
# ===========================================================================


_CONVERTERS: list[type[BaseConverter]] = [
    PDFConverter,
    DOCXConverter,
    PPTXConverter,
    XLSConverter,
    XLSXConverter,
    CSVConverter,
    HTMLConverter,
    TextConverter,
    ImageFileConverter,
]
_SUPPORTED_EXT: set[str] = (
    {e for cls in _CONVERTERS for e in cls.SUPPORTED_EXTENSIONS} | set(LEGACY_OFFICE_EXTENSIONS)
)


def get_converter(filename: str) -> BaseConverter:
    suffix = Path(filename).suffix.lower()
    for cls in _CONVERTERS:
        if suffix in cls.SUPPORTED_EXTENSIONS:
            return cls()
    raise ValueError(f"ไม่รองรับ '{suffix}'. รองรับ: {sorted(_SUPPORTED_EXT)}")


# ===========================================================================
# ── EMBEDDING HELPERS
# ===========================================================================


def _embed(emb_obj: Any, text: str) -> list[float] | None:
    if emb_obj is None:
        return None
    try:
        if hasattr(emb_obj, "embed_query"):
            return emb_obj.embed_query(text)
        if hasattr(emb_obj, "embed_documents"):
            r = emb_obj.embed_documents([text])
            return r[0] if r else None
    except Exception:
        return None
    return None


def _embed_with_error(emb_obj: Any, text: str) -> tuple[list[float] | None, str | None]:
    if emb_obj is None:
        return None, "embedding object is None"
    try:
        if hasattr(emb_obj, "embed_query"):
            vec = emb_obj.embed_query(text)
            return (vec, None) if vec else (None, "empty vector from embed_query")
        if hasattr(emb_obj, "embed_documents"):
            rows = emb_obj.embed_documents([text])
            if rows and rows[0]:
                return rows[0], None
            return None, "empty vector from embed_documents"
        return None, "embedding object has no embed_query/embed_documents"
    except Exception as e:
        return None, str(e)


def detect_dim(emb_obj: Any) -> int | None:
    if emb_obj is None:
        return None
    vec = _embed(emb_obj, "test")
    if vec and len(vec) > 0:
        return len(vec)
    return None


def resolve_vector_dim(
    embeddings: Any | None,
    *,
    configured_dim: int = 0,
    use_embedding: bool = True,
) -> int:
    if configured_dim and configured_dim > 0:
        return int(configured_dim)
    if use_embedding and embeddings is not None:
        detected = detect_dim(embeddings)
        if detected and detected > 0:
            return detected
    return 2048


def _embed_image_with_error(embeddings: Any, image_url: str) -> tuple[list[float] | None, str | None]:
    if embeddings is None:
        return None, "embedding object is None"
    try:
        if hasattr(embeddings, "embed_image_url"):
            vec = embeddings.embed_image_url(image_url)
            return (vec, None) if vec else (None, "empty vector from embed_image_url")
        return None, "embedding object does not support images"
    except Exception as e:
        return None, str(e)


def _insert_image_vector_rows(
    *,
    conn,
    vec_tbl: str,
    embeddings: Any,
    images: list[tuple[str, str]],
    base_chunk_index: int,
    doc_name: str | None,
    src: str,
    bucket_url: str,
    source_file_url: str | None,
    source_storage_path: str | None,
    actor_username: str | None,
    errors: list[str],
    progress_callback: Callable[[str, int, str | None], None] | None = None,
) -> list[Any]:
    """อัปโหลดรูปแล้ว embed แยกเป็นแถว vector (metadata.chunk_type=image)."""
    doc_ids: list[Any] = []
    total = len(images)
    if total == 0:
        return doc_ids

    for img_idx, (asset_name, pub_url) in enumerate(images, start=1):
        content = f"![{asset_name}]({pub_url})"
        chunk_index = base_chunk_index + img_idx
        try:
            vec, embed_err = _embed_image_with_error(embeddings, pub_url)
            if not vec:
                errors.append(f"image[{asset_name}]: embedding failed - {embed_err or 'unknown error'}")
                continue
            if progress_callback:
                progress_callback(
                    "embedding",
                    int((img_idx / total) * 100),
                    f"embedded image {img_idx}/{total}",
                )
            vec_str = f"[{','.join(str(float(x)) for x in vec)}]"
            with conn.cursor() as cur:
                cur.execute(
                    psql_sql.SQL(
                        """INSERT INTO {}
                       (content, document_name, chunk_index, token_count, created_by, updated_by, metadata, embedding)
                       VALUES (%s, %s, %s, %s, %s, %s, %s::jsonb, %s::vector) RETURNING id"""
                    ).format(psql_sql.Identifier(vec_tbl)),
                    (
                        content,
                        doc_name,
                        chunk_index,
                        0,
                        actor_username,
                        actor_username,
                        json.dumps(
                            {
                                "source_filename": src,
                                "document_name": doc_name,
                                "bucket_url": bucket_url,
                                "source_file_url": source_file_url,
                                "source_file_storage_path": source_storage_path,
                                "chunk_index": chunk_index,
                                "chunk_type": "image",
                                "image_url": pub_url,
                                "image_name": asset_name,
                                "token_count": 0,
                            },
                            ensure_ascii=False,
                        ),
                        vec_str,
                    ),
                )
                row = cur.fetchone()
                if row:
                    doc_ids.append(row[0])
            conn.commit()
            if progress_callback:
                progress_callback("db", int((img_idx / total) * 100), f"imported image {img_idx}/{total}")
        except Exception as e:
            conn.rollback()
            errors.append(f"image[{asset_name}]: {e}")
    return doc_ids


# ===========================================================================
# ── TEXT CHUNKER
# ===========================================================================


def chunk_text(text: str, size: int, overlap: int, split_text: str = "\n\n") -> list[str]:
    if not text.strip() or size <= 0:
        return [text] if text.strip() else []

    token_pattern = re.compile(r"\w+|[^\w\s]", flags=re.UNICODE)

    def _token_count(segment: str) -> int:
        if not segment:
            return 0
        return len(token_pattern.findall(segment))

    def _join_paragraphs(parts: list[str], sep: str) -> str:
        return sep.join(parts).strip()
    _ = overlap  # Overlap intentionally disabled by requirement.

    split_token = (split_text or "\n\n").replace("\\n", "\n")
    # Strict paragraph-only split: never cut inside paragraph/URL.
    paragraphs = [p.strip() for p in text.split(split_token) if p and p.strip()]
    if not paragraphs:
        paragraphs = [text.strip()]

    chunks: list[str] = []
    current_parts: list[str] = []
    sep = split_token or "\n\n"

    for para in paragraphs:
        if not current_parts:
            current_parts = [para]
            continue

        current_text = _join_paragraphs(current_parts, sep)
        current_tokens = _token_count(current_text)

        candidate_parts = current_parts + [para]
        candidate_text = _join_paragraphs(candidate_parts, sep)
        candidate_tokens = _token_count(candidate_text)

        # Choose split position that gets closest to CHUNK_TOKEN_SIZE.
        current_diff = abs(size - current_tokens)
        candidate_diff = abs(size - candidate_tokens)

        if candidate_tokens <= size or candidate_diff <= current_diff:
            current_parts = candidate_parts
            # Flush when already at/over target and best choice is current candidate.
            if candidate_tokens >= size:
                chunks.append(candidate_text)
                current_parts = []
            continue

        chunks.append(current_text)
        current_parts = [para]

    current_text = _join_paragraphs(current_parts, sep)
    if current_text:
        chunks.append(current_text)

    return chunks


def estimate_token_count(text: str) -> int:
    if not text:
        return 0
    # Lightweight tokenizer approximation aligned with chunk_text token slicing.
    return len(re.findall(r"\w+|[^\w\s]", text, flags=re.UNICODE))


def format_summary_progress_label(*, file_index: int, total_files: int, llm_text: str = "") -> str:
    return f"files {int(file_index)}/{int(total_files)} tokens {estimate_token_count(llm_text or '')}"


def make_summary_stream_emitter(
    progress_fn: Callable[[str, int, str | None], None],
    *,
    file_index: int,
    total_files: int,
) -> Callable[[str], None]:
    """Emit summary progress label whenever streamed LLM token count increases."""

    last_emitted_tokens = -1

    def emit(llm_text: str) -> None:
        nonlocal last_emitted_tokens
        tok = estimate_token_count(llm_text or "")
        if tok <= last_emitted_tokens:
            return
        last_emitted_tokens = tok
        label = format_summary_progress_label(
            file_index=file_index,
            total_files=total_files,
            llm_text=llm_text,
        )
        progress_fn("summary", min(99, max(1, tok)), label)

    return emit


def format_llm_chunks_progress_label(*, done: int, total: int, llm_text: str = "") -> str:
    tok = estimate_token_count(llm_text or "")
    if total <= 0:
        return f"chunks 1/1 tokens {tok}"
    return f"chunks {int(done)}/{int(total)} tokens {tok}"


def format_export_images_progress_label(*, done: int, total: int) -> str:
    """Label สำหรับ progress.export_images — ให้ UI แสดงนับรูปแบบ realtime เหมือน llm_chunks."""
    if total <= 0:
        return "images 0/0"
    return f"images {int(done)}/{int(total)}"


def _apply_image_url_map_to_markdown(markdown: str, url_map: dict[str, str]) -> str:
    """Replace relative image paths only for assets that were uploaded successfully."""
    body = markdown or ""
    if not url_map or not body:
        return body
    for old_path, public_url in url_map.items():
        body = body.replace(old_path, public_url)
    return body


def _upload_conversion_images(
    *,
    conv: ConversionResult,
    storage_dir: str,
    sb_client: Any,
    bucket: str,
    settings: Any,
    conn: Any,
    errors: list[str],
    progress_callback: Callable[[str, int, str | None], None] | None = None,
) -> tuple[list[str], list[tuple[str, str]], dict[str, str]]:
    """Upload extracted images; build URL map only for objects that landed in storage."""
    from app.document_api.core.storage_urls import public_storage_object_url

    def _progress(stage: str, percent: int, detail: str | None = None) -> None:
        if not progress_callback:
            return
        try:
            progress_callback(stage, max(0, min(100, int(percent))), detail)
        except Exception:
            pass

    image_urls: list[str] = []
    uploaded_images: list[tuple[str, str]] = []
    image_url_map: dict[str, str] = {}
    total_images = len(conv.images)
    if total_images == 0:
        _progress("images", 100, format_export_images_progress_label(done=0, total=0))
        return image_urls, uploaded_images, image_url_map

    _progress("images", 0, format_export_images_progress_label(done=0, total=total_images))
    for img_idx, asset in enumerate(conv.images, start=1):
        if not asset.path.exists():
            errors.append(f"image not found: {asset.path}")
            continue
        storage_path = f"{storage_dir}/{_safe_storage_name(asset.name)}"
        mime = mimetypes.guess_type(str(asset.path))[0] or "image/png"
        try:
            sb_client.storage.from_(bucket).upload(
                path=storage_path,
                file=asset.path.read_bytes(),
                file_options={"content-type": mime, "upsert": "true"},
            )
            pub_url = public_storage_object_url(bucket, storage_path, settings)
            image_urls.append(pub_url)
            uploaded_images.append((asset.name, pub_url))
            image_url_map[asset.rel_path] = pub_url
            image_url_map[str(asset.path)] = pub_url
            conn.commit()
            _progress(
                "images",
                int((img_idx / total_images) * 100),
                format_export_images_progress_label(done=img_idx, total=total_images),
            )
        except Exception as e:
            conn.rollback()
            errors.append(f"img {asset.name}: {e}")
    return image_urls, uploaded_images, image_url_map


def _image_upload_incomplete_error(*, uploaded: int, expected: int) -> str:
    return f"image upload incomplete: {uploaded}/{expected}"


def format_summary_progress_detail(text: str = "") -> str:
    return format_summary_progress_label(file_index=1, total_files=1, llm_text=text)


def format_llm_chunks_progress_detail(*, done: int, total: int, text: str = "") -> str:
    return format_llm_chunks_progress_label(done=done, total=total, llm_text=text)


class LlmChunksProgressHandles:
    """Callbacks for LLM rearrange progress — stream updates token count in real time."""

    __slots__ = (
        "chunk_progress_callback",
        "chunk_stream_callback",
        "start_heartbeat",
        "stop_heartbeat",
        "last_detail",
    )

    def __init__(
        self,
        progress_fn: Callable[[str, int, str | None], None],
        *,
        stream_enabled: bool = True,
    ) -> None:
        from app.document_api.services.rearrange_progress import RearrangeProgressHeartbeat

        _last_llm_chunks_detail = ""
        _last_emitted_tokens = -1
        _stream_chunk_cur = 0
        _sheet_export_stream_seen = False
        _rearrange_hb: RearrangeProgressHeartbeat | None = None

        def _emit_progress(pct: int, detail: str) -> None:
            nonlocal _last_llm_chunks_detail
            if not detail.startswith("chunks "):
                return
            _last_llm_chunks_detail = detail
            if _rearrange_hb is not None:
                _rearrange_hb.note_stream_pct(pct)
            progress_fn("llm_chunks", pct, detail)

        def _pct_for_stream(cur: int, total: int, acc: str, phase: str) -> int:
            acc_len = len(acc)
            denom = max(1, total)
            frac_within = min(1.0, 1.0 - math.exp(-acc_len / 1800.0))
            if phase == "sheet_export":
                raw_pct = frac_within * 35.0
            elif _sheet_export_stream_seen:
                raw_pct = 35.0 + ((cur - 1 + frac_within) / denom) * 65.0
            else:
                raw_pct = ((cur - 1 + frac_within) / denom) * 100.0
            return min(99, int(round(raw_pct)))

        def _chunk_prog(done: int, total: int, snippet: str = "", phase: str = "reflow_chunk") -> None:
            nonlocal _last_emitted_tokens, _stream_chunk_cur
            _stream_chunk_cur = done
            _last_emitted_tokens = -1
            if total <= 0:
                detail = format_llm_chunks_progress_label(done=1, total=1, llm_text=snippet)
                _emit_progress(100, detail)
                return
            if _sheet_export_stream_seen:
                raw_pct = 35.0 + (done / max(1, total)) * 65.0
            else:
                raw_pct = (done / max(1, total)) * 100.0
            pct = int(max(0, min(100, round(raw_pct))))
            detail = format_llm_chunks_progress_label(done=done, total=total, llm_text=snippet)
            _last_emitted_tokens = estimate_token_count(snippet or "")
            _emit_progress(pct, detail)

        def _chunk_stream(cur: int, total: int, acc: str, phase: str = "reflow_chunk") -> None:
            nonlocal _last_emitted_tokens, _stream_chunk_cur, _sheet_export_stream_seen
            if phase == "sheet_export":
                _sheet_export_stream_seen = True
            if cur != _stream_chunk_cur:
                _stream_chunk_cur = cur
                _last_emitted_tokens = -1
            detail = format_llm_chunks_progress_label(done=cur, total=total, llm_text=acc)
            tok = estimate_token_count(acc or "")
            if tok <= _last_emitted_tokens:
                return
            _last_emitted_tokens = tok
            _emit_progress(_pct_for_stream(cur, total, acc, phase), detail)

        def _start_heartbeat(phase: str = "rearrange") -> None:
            nonlocal _rearrange_hb

            def _pulse(pct: int, detail: str) -> None:
                if _last_llm_chunks_detail.startswith("chunks "):
                    progress_fn("llm_chunks", pct, _last_llm_chunks_detail)
                else:
                    progress_fn("llm_chunks", pct, detail)

            _rearrange_hb = RearrangeProgressHeartbeat(_pulse, interval_sec=2.0)
            _rearrange_hb.start(phase)

        def _stop_heartbeat() -> None:
            nonlocal _rearrange_hb
            if _rearrange_hb is not None:
                _rearrange_hb.stop()
                _rearrange_hb = None

        self.chunk_progress_callback = _chunk_prog
        self.chunk_stream_callback = _chunk_stream if stream_enabled else None
        self.start_heartbeat = _start_heartbeat
        self.stop_heartbeat = _stop_heartbeat
        self.last_detail = lambda: _last_llm_chunks_detail


# ===========================================================================
# ── SUPABASE DDL SETUP (statement-by-statement; same logic as your v5)
# ===========================================================================


def _vector_chunks_setup_statements(table: str, dim: int) -> list[tuple[str, str]]:
    """DDL สำหรับตารางเวกเตอร์ + ฟังก์ชัน match_documents (`table` ต้องผ่าน _vector_table_ident แล้ว)."""
    t = table
    d = max(1, int(dim))
    emb_idx = f"{t}_embedding_idx"[:63]
    out: list[tuple[str, str]] = [
        ("extension_vector", "CREATE EXTENSION IF NOT EXISTS vector"),
        ("drop_fn_old_1", "DROP FUNCTION IF EXISTS match_documents(vector, int, jsonb)"),
        ("drop_fn_old_2", "DROP FUNCTION IF EXISTS match_documents(vector(768), int, jsonb)"),
        ("drop_fn_old_3", "DROP FUNCTION IF EXISTS match_documents(vector(1024), int, jsonb)"),
        ("drop_fn_old_4", "DROP FUNCTION IF EXISTS match_documents(vector(2048), int, jsonb)"),
    ]
    # อย่า DROP ตาราง legacy ชื่อ documents ถ้าผู้ใช้กำลังใช้ชื่อนั้นเป็นตารางเวกเตอร์
    if t != "documents":
        out.append(("drop_legacy_documents_table", "DROP TABLE IF EXISTS documents"))
    out.extend(
        [
            (
                "table_vector_chunks",
                f"""CREATE TABLE IF NOT EXISTS {t} (
         id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
         content TEXT,
         document_name TEXT,
         chunk_index INT,
         token_count INT,
         created_by TEXT,
         updated_by TEXT,
         created_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
         updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
         metadata JSONB,
         embedding vector({d})
     )""",
            ),
            (
                "alter_vector_chunks_document_name",
                f"ALTER TABLE {t} ADD COLUMN IF NOT EXISTS document_name TEXT",
            ),
            (
                "alter_vector_chunks_chunk_index",
                f"ALTER TABLE {t} ADD COLUMN IF NOT EXISTS chunk_index INT",
            ),
            (
                "alter_vector_chunks_token_count",
                f"ALTER TABLE {t} ADD COLUMN IF NOT EXISTS token_count INT",
            ),
            (
                "alter_vector_chunks_created_by",
                f"ALTER TABLE {t} ADD COLUMN IF NOT EXISTS created_by TEXT",
            ),
            (
                "alter_vector_chunks_updated_by",
                f"ALTER TABLE {t} ADD COLUMN IF NOT EXISTS updated_by TEXT",
            ),
            (
                "alter_vector_chunks_created_at",
                f"ALTER TABLE {t} ADD COLUMN IF NOT EXISTS created_at TIMESTAMPTZ NOT NULL DEFAULT NOW()",
            ),
            (
                "alter_vector_chunks_updated_at",
                f"ALTER TABLE {t} ADD COLUMN IF NOT EXISTS updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()",
            ),
            (
                "index_embedding",
                f"""CREATE INDEX IF NOT EXISTS {emb_idx}
        ON {t} USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100)""",
            ),
            (
                "fn_match_documents",
                f"""CREATE OR REPLACE FUNCTION match_documents(
         query_embedding vector({d}), match_count INT DEFAULT 4,
         filter JSONB DEFAULT '{{}}'
     ) RETURNS TABLE (id UUID, content TEXT, metadata JSONB, similarity FLOAT)
     LANGUAGE plpgsql AS $$
     BEGIN
         RETURN QUERY EXECUTE format(
         'SELECT d.id, d.content, d.metadata,
                1 - (d.embedding <=> $1) AS similarity
         FROM %I d
         WHERE d.embedding IS NOT NULL
           AND d.metadata @> $2::jsonb
         ORDER BY d.embedding <=> $1
         LIMIT $3',
         '{t}'
         ) USING query_embedding, filter, match_count;
     END; $$""",
            ),
        ]
    )
    return out


_SETUP_STATEMENTS_TAIL: list[tuple[str, str]] = [
    (
        "table_transcript",
        """CREATE TABLE IF NOT EXISTS transcript (
         id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
         document_name TEXT NOT NULL,
         transcript_name TEXT NOT NULL DEFAULT '',
         content TEXT,
         files JSONB NOT NULL DEFAULT '[]'::jsonb,
         segments JSONB NOT NULL DEFAULT '[]'::jsonb,
         created_by TEXT,
         updated_by TEXT,
         created_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
         updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
     )""",
    ),
    (
        "index_transcript_document_name",
        "CREATE INDEX IF NOT EXISTS transcript_document_name_idx ON transcript (document_name)",
    ),
    (
        "alter_transcript_segments",
        "ALTER TABLE transcript ADD COLUMN IF NOT EXISTS segments JSONB NOT NULL DEFAULT '[]'::jsonb",
    ),
    (
        "alter_transcript_transcript_name",
        "ALTER TABLE transcript ADD COLUMN IF NOT EXISTS transcript_name TEXT NOT NULL DEFAULT ''",
    ),
    (
        "index_transcript_document_transcript_name",
        "CREATE INDEX IF NOT EXISTS transcript_document_transcript_name_idx ON transcript (document_name, transcript_name)",
    ),
    (
        "alter_transcript_audio_llm_summary",
        "ALTER TABLE transcript ADD COLUMN IF NOT EXISTS audio_llm_summary TEXT",
    ),
    (
        "alter_transcript_audio_llm_report",
        "ALTER TABLE transcript ADD COLUMN IF NOT EXISTS audio_llm_report TEXT",
    ),
    (
        "table_document_ingest_pending",
        """CREATE TABLE IF NOT EXISTS document_ingest_pending (
         id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
         document_name TEXT NOT NULL,
         source_filename TEXT NOT NULL,
         markdown_text TEXT NOT NULL,
         bucket_url TEXT,
         source_file_storage_path TEXT,
         source_file_url TEXT,
         converter_metadata JSONB,
         rearrange_notes JSONB,
         rearrange_llm_raw TEXT,
         llm_summary TEXT,
         chunk_size INT NOT NULL DEFAULT 1500,
         chunk_overlap INT NOT NULL DEFAULT 200,
         split_text TEXT NOT NULL DEFAULT E'\\n\\n',
         on_duplicate TEXT NOT NULL DEFAULT 'replace',
         status TEXT NOT NULL DEFAULT 'pending',
         job_id TEXT,
         created_by TEXT,
         updated_by TEXT,
         created_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
         updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
     )""",
    ),
    (
        "migrate_document_ingest_pending_llm_summary_rename",
        """DO $m$
        BEGIN
          IF EXISTS (
            SELECT 1 FROM information_schema.columns
            WHERE table_schema = 'public' AND table_name = 'document_ingest_pending'
              AND column_name = 'admin_llm_summary'
          ) AND NOT EXISTS (
            SELECT 1 FROM information_schema.columns
            WHERE table_schema = 'public' AND table_name = 'document_ingest_pending'
              AND column_name = 'llm_summary'
          ) THEN
            ALTER TABLE public.document_ingest_pending
              RENAME COLUMN admin_llm_summary TO llm_summary;
          END IF;
        END $m$""",
    ),
    (
        "index_document_ingest_pending_doc_status",
        "CREATE INDEX IF NOT EXISTS document_ingest_pending_doc_status_idx ON document_ingest_pending (document_name, status)",
    ),
    (
        "index_document_ingest_pending_job",
        "CREATE INDEX IF NOT EXISTS document_ingest_pending_job_idx ON document_ingest_pending (job_id)",
    ),
    (
        "alter_document_ingest_pending_summary_ready",
        "ALTER TABLE document_ingest_pending ADD COLUMN IF NOT EXISTS summary_ready BOOLEAN NOT NULL DEFAULT false",
    ),
    (
        "backfill_document_ingest_pending_summary_ready",
        """UPDATE document_ingest_pending
           SET summary_ready = true
           WHERE status = 'pending'
             AND summary_ready = false""",
    ),
    (
        "alter_document_folder_approved_by",
        "ALTER TABLE document_folder ADD COLUMN IF NOT EXISTS approved_by TEXT",
    ),
    (
        "alter_document_folder_approved_at",
        "ALTER TABLE document_folder ADD COLUMN IF NOT EXISTS approved_at TIMESTAMPTZ",
    ),
    (
        "backfill_document_folder_approval_audit",
        """UPDATE document_folder
           SET approved_by = updated_by,
               approved_at = updated_at
           WHERE approved_by IS NULL
             AND updated_by IS NOT NULL""",
    ),
]


def _build_statements(dim: int, table_name: str) -> list[tuple[str, str]]:
    t = _vector_table_ident(table_name)
    return _vector_chunks_setup_statements(t, dim) + _SETUP_STATEMENTS_TAIL


def _get_existing_dim(conn, table_name: str) -> int | None:
    ident = _vector_table_ident(table_name)
    try:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT a.atttypmod
                FROM   pg_attribute a
                JOIN   pg_class     c ON c.oid = a.attrelid
                WHERE  c.relname = %s
                  AND  a.attname = 'embedding'
                  AND  a.attnum  > 0
                  AND  NOT a.attisdropped
            """,
                (ident,),
            )
            row = cur.fetchone()
            if row and row[0] and row[0] > 0:
                return int(row[0])
    except Exception:
        conn.rollback()
    return None


def ensure_schema(conn) -> list[str]:
    _ = conn
    return []


def setup_ddl(conn, dim: int, table_name: str | None = None) -> None:
    from app.document_api.core.config import get_settings

    raw = table_name if table_name is not None else get_settings().supabase_table_name
    ident = _vector_table_ident(raw)
    _ = _get_existing_dim(conn, ident)

    for _, sql in _build_statements(dim, ident):
        s = sql.strip()
        if not s:
            continue
        try:
            with conn.cursor() as cur:
                cur.execute(s)
            conn.commit()
        except Exception:
            conn.rollback()


# ===========================================================================
# ── PIPELINE (convert / upload / query)
# ===========================================================================


@dataclass
class PgConfig:
    host: str = "127.0.0.1"
    port: int = 5432
    database: str = "postgres"
    user: str = "postgres"
    password: str = ""
    sslmode: str = "disable"


@dataclass
class SupabaseConfig:
    url: str = ""
    service_key: str = ""
    storage_bucket: str = "document-images"
    transcript_bucket: str = "transcript"
    table_name: str = "langflow"
    query_name: str = "match_documents"


def _pg_conn(cfg: PgConfig):
    import psycopg2

    return psycopg2.connect(
        host=cfg.host,
        port=cfg.port,
        database=cfg.database,
        user=cfg.user,
        password=cfg.password,
        sslmode=cfg.sslmode,
    )


def _sb_client(cfg: SupabaseConfig):
    from supabase.client import create_client

    return create_client(cfg.url, supabase_key=cfg.service_key)


def download_ingest_source_bytes(document_name: str, source_filename: str) -> bytes:
    """Download original source file from Supabase document storage (for job retry)."""
    from app.document_api.core.config import get_settings

    s = get_settings()
    if not (s.supabase_url or "").strip() or not (s.supabase_service_key or "").strip():
        raise ValueError("Supabase is not configured")
    doc = (document_name or "").strip()
    src = (source_filename or "").strip()
    if not src:
        raise ValueError("source_filename is required")
    storage_root = _safe_storage_name(doc or src)
    path = f"{storage_root}/source/{_safe_storage_name(src)}"
    sb = SupabaseConfig(
        url=s.supabase_url,
        service_key=s.supabase_service_key,
        storage_bucket=s.supabase_storage_bucket,
        transcript_bucket=s.supabase_transcript_bucket,
        table_name=s.supabase_table_name,
        query_name=s.supabase_query_name,
    )
    client = _sb_client(sb)
    try:
        data = client.storage.from_(s.supabase_storage_bucket).download(path)
    except Exception as e:
        raise ValueError(f"storage download failed for {path}: {e}") from e
    if not data:
        raise ValueError(f"storage download returned empty payload for {path}")
    return bytes(data)


def _ensure_bucket(sb, bucket: str):
    try:
        names = [b.name for b in sb.storage.list_buckets()]
        if bucket not in names:
            sb.storage.create_bucket(bucket, options={"public": True})
    except Exception:
        pass


def bootstrap_on_startup() -> None:
    """เรียกตอนสตาร์ทแอป: DDL Postgres, ตาราง document_folder, สร้าง storage bucket หลัก."""
    import logging

    from app.document_api.core.config import get_settings
    from app.document_api.services.folder_catalog import ensure_folder_table

    log = logging.getLogger(__name__)
    settings = get_settings()

    pg = PgConfig(
        host=settings.pg_host,
        port=settings.pg_port,
        database=settings.pg_database,
        user=settings.pg_user,
        password=settings.pg_password,
        sslmode=settings.pg_sslmode,
    )
    dim = max(1, int(settings.embedding_dimensions or 1024))

    try:
        conn = _pg_conn(pg)
        try:
            setup_ddl(conn, dim)
            ensure_schema(conn)
        finally:
            conn.close()
    except Exception as e:
        log.warning("startup: database DDL skipped (%s)", e)

    try:
        ensure_folder_table()
    except Exception as e:
        log.warning("startup: document_folder init skipped (%s)", e)

    if not (settings.supabase_url or "").strip() or not (settings.supabase_service_key or "").strip():
        return

    sb = SupabaseConfig(
        url=settings.supabase_url,
        service_key=settings.supabase_service_key,
        storage_bucket=settings.supabase_storage_bucket,
        transcript_bucket=settings.supabase_transcript_bucket,
        table_name=settings.supabase_table_name,
        query_name=settings.supabase_query_name,
    )
    try:
        client = _sb_client(sb)
        _ensure_bucket(client, sb.storage_bucket)
        _ensure_bucket(client, sb.transcript_bucket)
    except Exception as e:
        log.warning("startup: storage buckets init skipped (%s)", e)


def convert_file_to_markdown(file_path: Path) -> ConversionResult:
    import logging

    from app.document_api.services.office_legacy_convert import (
        OfficeLegacyConvertError,
        convert_legacy_office_file,
        is_legacy_office_extension,
    )

    log = logging.getLogger(__name__)
    out_dir = Path(tempfile.mkdtemp(prefix="document_api_conv_"))
    src = file_path
    suffix = file_path.suffix.lower()
    if is_legacy_office_extension(suffix):
        try:
            src = convert_legacy_office_file(file_path, out_dir)
        except OfficeLegacyConvertError as exc:
            if suffix == ".xls" and XLSConverter.is_available():
                log.warning("LibreOffice unavailable for %s; using xlrd fallback", file_path.name)
                src = file_path
            else:
                raise ValueError(str(exc)) from exc
    return get_converter(src.name).convert(src, out_dir)


def ingest_to_supabase(
    *,
    conv: ConversionResult,
    enable_supabase: bool,
    pg: PgConfig,
    sb: SupabaseConfig,
    on_duplicate: str = "replace",
    chunk_size: int = 1500,
    chunk_overlap: int = 200,
    split_text: str = "\n\n",
    document_name: str | None = None,
    source_file_bytes: bytes | None = None,
    embeddings: Any | None = None,
    progress_callback: Callable[[str, int, str | None], None] | None = None,
    rearrange_stream_callback: Callable[[str, str], None] | None = None,
    actor_username: str | None = None,
    department_id: str | None = None,
    defer_vector_commit: bool = False,
    job_id: str | None = None,
) -> dict:
    def _progress(stage: str, percent: int, detail: str | None = None) -> None:
        if not progress_callback:
            return
        try:
            progress_callback(stage, max(0, min(100, int(percent))), detail)
        except Exception:
            pass

    if not enable_supabase:
        _progress("llm_chunks", 100, "llm rearrange skipped")
        _progress("images", 100, format_export_images_progress_label(done=0, total=0))
        _progress("embedding", 100, "embedding skipped")
        _progress("db", 100, "database import skipped")
        return {
            "status": "disabled",
            "source_filename": conv.source_filename,
            "doc_ids": [],
            "image_urls": [],
            "errors": [],
            "toc_errors": [],
            "rearrange_notes": [],
            "rearrange_llm_raw": "",
            "rearrange_llm_raw_chars": 0,
            "chunks_uploaded": 0,
            "images_uploaded": 0,
            "embedding_used": False,
            "vector_dim": None,
            "pending_ingest_id": None,
        }

    conn = _pg_conn(pg)
    sb_client = _sb_client(sb)

    use_emb = embeddings is not None
    from app.document_api.core.config import get_settings as _get_settings
    from app.document_api.core.storage_urls import public_storage_bucket_base

    _emb_settings = _get_settings()
    dim = resolve_vector_dim(
        embeddings,
        configured_dim=_emb_settings.embedding_dimensions,
        use_embedding=use_emb,
    )

    setup_ddl(conn, dim, sb.table_name)
    ensure_schema(conn)
    vec_tbl = _vector_table_ident(sb.table_name)
    _ensure_bucket(sb_client, sb.storage_bucket)

    src = conv.source_filename
    doc_name = (document_name or "").strip() or None
    storage_root = _safe_storage_name(doc_name or src)
    storage_dir = f"{storage_root}/files"
    source_storage_path = f"{storage_root}/source/{_safe_storage_name(src)}"
    bucket_url = public_storage_bucket_base(sb.storage_bucket, _emb_settings)
    source_file_url: str | None = None
    errors: list[str] = []

    if source_file_bytes:
        src_mime = mimetypes.guess_type(src)[0] or "application/octet-stream"
        try:
            sb_client.storage.from_(sb.storage_bucket).upload(
                path=source_storage_path,
                file=source_file_bytes,
                file_options={"content-type": src_mime, "upsert": "true"},
            )
            source_file_url = f"{bucket_url}/{source_storage_path}"
        except Exception as e:
            errors.append(f"source_file {src}: {e}")

    if on_duplicate == "replace" and not defer_vector_commit:
        try:
            where_sql, scope_params = _vector_rows_for_source_sql(doc_name)
            with conn.cursor() as cur:
                cur.execute(
                    psql_sql.SQL("DELETE FROM {} WHERE ").format(psql_sql.Identifier(vec_tbl))
                    + psql_sql.SQL(where_sql),
                    (src, *scope_params),
                )
            conn.commit()
        except Exception:
            conn.rollback()
    elif on_duplicate == "skip":
        try:
            where_sql, scope_params = _vector_rows_for_source_sql(doc_name)
            with conn.cursor() as cur:
                cur.execute(
                    psql_sql.SQL("SELECT 1 FROM {} WHERE ").format(psql_sql.Identifier(vec_tbl))
                    + psql_sql.SQL(where_sql)
                    + psql_sql.SQL(" LIMIT 1"),
                    (src, *scope_params),
                )
                if cur.fetchone():
                    conn.close()
                    return {
                        "status": "skipped",
                        "source_filename": src,
                        "doc_ids": [],
                        "image_urls": [],
                        "errors": [],
                        "toc_errors": [],
                        "rearrange_notes": [],
                        "rearrange_llm_raw": "",
                        "rearrange_llm_raw_chars": 0,
                        "chunks_uploaded": 0,
                        "images_uploaded": 0,
                        "embedding_used": use_emb,
                        "vector_dim": dim,
                    }
        except Exception:
            conn.rollback()

    doc_ids: list[int] = []
    image_urls: list[str] = []
    uploaded_images: list[tuple[str, str]] = []
    image_url_map: dict[str, str] = {}
    images_expected = len(conv.images)

    md_body = conv.markdown or ""
    rearrange_notes: list[str] = []
    rearrange_llm_raw = ""

    if defer_vector_commit:
        image_urls, uploaded_images, image_url_map = _upload_conversion_images(
            conv=conv,
            storage_dir=storage_dir,
            sb_client=sb_client,
            bucket=sb.storage_bucket,
            settings=_emb_settings,
            conn=conn,
            errors=errors,
            progress_callback=progress_callback,
        )
        md_body = _apply_image_url_map_to_markdown(md_body, image_url_map)
        if images_expected > 0 and len(image_urls) < images_expected:
            errors.append(_image_upload_incomplete_error(uploaded=len(image_urls), expected=images_expected))

    if defer_vector_commit:
        from app.document_api.services.pending_ingest_catalog import replace_pending_ingest

        pending_id = replace_pending_ingest(
            document_name=doc_name or "",
            source_filename=src,
            markdown_text=md_body,
            bucket_url=bucket_url,
            source_file_storage_path=source_storage_path,
            source_file_url=source_file_url,
            converter_metadata=dict(conv.metadata or {}),
            rearrange_notes=rearrange_notes,
            rearrange_llm_raw=rearrange_llm_raw or "",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            split_text=split_text,
            on_duplicate=on_duplicate,
            job_id=job_id,
            created_by=actor_username,
            department_id=department_id,
        )
        conn.close()
        return {
            "status": "pending_approval",
            "pending_ingest_id": pending_id,
            "source_filename": src,
            "doc_ids": [],
            "chunks_uploaded": 0,
            "images_uploaded": len(image_urls),
            "images_expected": images_expected,
            "image_urls": image_urls,
            "image_url_map": image_url_map,
            "bucket_url": bucket_url,
            "source_file_url": source_file_url,
            "source_file_storage_path": source_storage_path,
            "errors": errors,
            "toc_errors": [],
            "rearrange_notes": rearrange_notes,
            "rearrange_llm_raw": rearrange_llm_raw or "",
            "rearrange_llm_raw_chars": len(rearrange_llm_raw or ""),
            "embedding_used": False,
            "vector_dim": dim,
            "table_name": sb.table_name,
            "query_name": sb.query_name,
            "llm_summary": "",
        }

    if not defer_vector_commit:
        image_urls, uploaded_images, image_url_map = _upload_conversion_images(
            conv=conv,
            storage_dir=storage_dir,
            sb_client=sb_client,
            bucket=sb.storage_bucket,
            settings=_emb_settings,
            conn=conn,
            errors=errors,
            progress_callback=progress_callback,
        )
        md_body = _apply_image_url_map_to_markdown(md_body, image_url_map)
        if images_expected > 0 and len(image_urls) < images_expected:
            errors.append(_image_upload_incomplete_error(uploaded=len(image_urls), expected=images_expected))
            conn.close()
            return {
                "status": "error",
                "source_filename": src,
                "doc_ids": [],
                "chunks_uploaded": 0,
                "images_uploaded": len(image_urls),
                "images_expected": images_expected,
                "image_urls": image_urls,
                "image_url_map": image_url_map,
                "bucket_url": bucket_url,
                "source_file_url": source_file_url,
                "source_file_storage_path": source_storage_path,
                "errors": errors,
                "toc_errors": [],
                "rearrange_notes": rearrange_notes,
                "rearrange_llm_raw": rearrange_llm_raw or "",
                "rearrange_llm_raw_chars": len(rearrange_llm_raw or ""),
                "embedding_used": use_emb,
                "vector_dim": dim,
                "table_name": sb.table_name,
                "query_name": sb.query_name,
                "pending_ingest_id": None,
            }
        _progress("llm_chunks", 100, format_llm_chunks_progress_label(done=0, total=0, llm_text=""))

    chunks = chunk_text(md_body, chunk_size, chunk_overlap, split_text=split_text) or (
        [md_body] if md_body else []
    )
    total_chunks = len(chunks)
    if total_chunks == 0:
        _progress("embedding", 100, "no chunks to embed")
        _progress("db", 100, "no chunks to import")

    if chunks and use_emb:
        for idx, chunk in enumerate(chunks):
            content_hash = hashlib.md5(chunk.encode()).hexdigest()
            token_count = estimate_token_count(chunk)
            try:
                vec, embed_err = _embed_with_error(embeddings, chunk)
                if not vec:
                    errors.append(f"doc[{idx}]: embedding failed - {embed_err or 'unknown error'}")
                    continue
                _progress("embedding", int(((idx + 1) / total_chunks) * 100), f"embedded chunk {idx + 1}/{total_chunks}")
                vec_str = f"[{','.join(str(float(x)) for x in vec)}]"
                with conn.cursor() as cur:
                    cur.execute(
                        psql_sql.SQL(
                            """INSERT INTO {}
                           (content, document_name, chunk_index, token_count, created_by, updated_by, metadata, embedding)
                           VALUES (%s, %s, %s, %s, %s, %s, %s::jsonb, %s::vector) RETURNING id"""
                        ).format(psql_sql.Identifier(vec_tbl)),
                        (
                            chunk,
                            doc_name,
                            idx,
                            token_count,
                            actor_username,
                            actor_username,
                            json.dumps(
                                {
                                    "source_filename": src,
                                    "document_name": doc_name,
                                    "bucket_url": bucket_url,
                                    "source_file_url": source_file_url,
                                    "source_file_storage_path": source_storage_path,
                                    "chunk_index": idx,
                                    "chunk_type": "text",
                                    "token_count": token_count,
                                    "content_hash": content_hash,
                                },
                                ensure_ascii=False,
                            ),
                            vec_str,
                        ),
                    )
                    row = cur.fetchone()
                    if row:
                        doc_ids.append(row[0])
                conn.commit()
                _progress("db", int(((idx + 1) / total_chunks) * 100), f"imported chunk {idx + 1}/{total_chunks}")
            except Exception as e:
                conn.rollback()
                errors.append(f"doc[{idx}]: {e}")

    elif chunks and not use_emb:
        _progress("embedding", 100, "embedding disabled")
        for idx, chunk in enumerate(chunks):
            content_hash = hashlib.md5(chunk.encode()).hexdigest()
            token_count = estimate_token_count(chunk)
            try:
                with conn.cursor() as cur:
                    cur.execute(
                        psql_sql.SQL(
                            """INSERT INTO {}
                           (content, document_name, chunk_index, token_count, created_by, updated_by, metadata)
                           VALUES (%s, %s, %s, %s, %s, %s, %s::jsonb) RETURNING id"""
                        ).format(psql_sql.Identifier(vec_tbl)),
                        (
                            chunk,
                            doc_name,
                            idx,
                            token_count,
                            actor_username,
                            actor_username,
                            json.dumps(
                                {
                                    "source_filename": src,
                                    "document_name": doc_name,
                                    "bucket_url": bucket_url,
                                    "source_file_url": source_file_url,
                                    "source_file_storage_path": source_storage_path,
                                    "chunk_index": idx,
                                    "chunk_type": "text",
                                    "token_count": token_count,
                                    "content_hash": content_hash,
                                },
                                ensure_ascii=False,
                            ),
                        ),
                    )
                    row = cur.fetchone()
                    if row:
                        doc_ids.append(row[0])
                conn.commit()
                _progress("db", int(((idx + 1) / total_chunks) * 100), f"imported chunk {idx + 1}/{total_chunks}")
            except Exception as e:
                conn.rollback()
                errors.append(f"doc[{idx}]: {e}")

    if use_emb and uploaded_images:
        img_doc_ids = _insert_image_vector_rows(
            conn=conn,
            vec_tbl=vec_tbl,
            embeddings=embeddings,
            images=uploaded_images,
            base_chunk_index=total_chunks,
            doc_name=doc_name,
            src=src,
            bucket_url=bucket_url,
            source_file_url=source_file_url,
            source_storage_path=source_storage_path,
            actor_username=actor_username,
            errors=errors,
            progress_callback=_progress,
        )
        doc_ids.extend(img_doc_ids)

    conn.close()

    toc_errors: list[str] = []

    return {
        "status": "success" if not errors else "partial",
        "source_filename": src,
        "doc_ids": doc_ids,
        "chunks_uploaded": len(doc_ids),
        "images_uploaded": len(image_urls),
        "images_expected": images_expected,
        "image_urls": image_urls,
        "image_url_map": image_url_map,
        "bucket_url": bucket_url,
        "source_file_url": source_file_url,
        "source_file_storage_path": source_storage_path,
        "errors": errors,
        "toc_errors": toc_errors,
        "rearrange_notes": rearrange_notes,
        "rearrange_llm_raw": rearrange_llm_raw or "",
        "rearrange_llm_raw_chars": len(rearrange_llm_raw or ""),
        "embedding_used": use_emb,
        "vector_dim": dim,
        "table_name": sb.table_name,
        "query_name": sb.query_name,
        "pending_ingest_id": None,
    }


def commit_pending_ingest_to_supabase(
    *,
    pending_id: str,
    pg: PgConfig,
    sb: SupabaseConfig,
    embeddings: Any | None = None,
    progress_callback: Callable[[str, int, str | None], None] | None = None,
    actor_username: str | None = None,
) -> dict:
    """นำแถว document_ingest_pending ไปแบ่ง chunk + embedding + TOC + ลง document_folder แล้วลบ pending"""
    from app.document_api.services.pending_ingest_catalog import delete_pending_row, get_pending_by_id

    row = get_pending_by_id(pending_id)
    if not row:
        return {"status": "error", "source_filename": "", "errors": ["ไม่พบ pending ingest"]}
    if row.get("status") != "pending":
        return {
            "status": "error",
            "source_filename": row.get("source_filename") or "",
            "errors": [f"สถานะไม่ใช่ pending: {row.get('status')}"],
        }

    def _progress(stage: str, percent: int, detail: str | None = None) -> None:
        if not progress_callback:
            return
        try:
            progress_callback(stage, max(0, min(100, int(percent))), detail)
        except Exception:
            pass

    md_body = row["markdown_text"] or ""
    src = row["source_filename"]
    doc_name = (row["document_name"] or "").strip() or None
    chunk_size = int(row["chunk_size"])
    chunk_overlap = int(row["chunk_overlap"])
    split_text = row["split_text"] or "\n\n"
    on_duplicate = row["on_duplicate"] or "replace"
    bucket_url = row["bucket_url"] or ""
    source_file_url = row["source_file_url"]
    source_storage_path = row["source_file_storage_path"]

    conn = _pg_conn(pg)
    sb_client = _sb_client(sb)
    use_emb = embeddings is not None
    from app.document_api.core.config import get_settings as _get_settings

    _emb_settings = _get_settings()
    dim = resolve_vector_dim(
        embeddings,
        configured_dim=_emb_settings.embedding_dimensions,
        use_embedding=use_emb,
    )

    setup_ddl(conn, dim, sb.table_name)
    ensure_schema(conn)
    vec_tbl = _vector_table_ident(sb.table_name)
    _ensure_bucket(sb_client, sb.storage_bucket)

    errors: list[str] = []
    doc_ids: list[int] = []
    image_urls: list[str] = []
    uploaded_images: list[tuple[str, str]] = []

    if on_duplicate == "replace":
        try:
            where_sql, scope_params = _vector_rows_for_source_sql(doc_name)
            with conn.cursor() as cur:
                cur.execute(
                    psql_sql.SQL("DELETE FROM {} WHERE ").format(psql_sql.Identifier(vec_tbl))
                    + psql_sql.SQL(where_sql),
                    (src, *scope_params),
                )
            conn.commit()
        except Exception:
            conn.rollback()

    storage_root = _safe_storage_name(doc_name or src)
    storage_dir = f"{storage_root}/files"
    conv = None
    try:
        raw_bytes = download_ingest_source_bytes(doc_name or "", src)
        tmp_dir = Path(tempfile.mkdtemp())
        try:
            tmp_path = tmp_dir / Path(src).name
            tmp_path.write_bytes(raw_bytes)
            conv = convert_file_to_markdown(tmp_path)
        finally:
            import shutil

            shutil.rmtree(tmp_dir, ignore_errors=True)
    except Exception as e:
        # Stored markdown is enough to commit; reconvert only refreshes images.
        if not md_body.strip():
            errors.append(f"reconvert: {e}")
        else:
            import logging

            logging.getLogger(__name__).warning(
                "commit_pending: reconvert skipped for %s (%s); using stored markdown",
                src,
                e,
            )

    rearrange_notes: list[str] = []
    rearrange_llm_raw = ""
    _progress("llm_chunks", 100, format_llm_chunks_progress_label(done=0, total=0, llm_text=""))

    images_expected = len(conv.images) if conv else 0
    image_url_map: dict[str, str] = {}
    if conv:
        image_urls, uploaded_images, image_url_map = _upload_conversion_images(
            conv=conv,
            storage_dir=storage_dir,
            sb_client=sb_client,
            bucket=sb.storage_bucket,
            settings=_emb_settings,
            conn=conn,
            errors=errors,
            progress_callback=progress_callback,
        )
        if image_url_map:
            md_body = _apply_image_url_map_to_markdown(md_body, image_url_map)
        if images_expected > 0 and len(image_urls) < images_expected:
            errors.append(_image_upload_incomplete_error(uploaded=len(image_urls), expected=images_expected))
            conn.close()
            return {
                "status": "error",
                "pending_ingest_id": pending_id,
                "source_filename": src,
                "doc_ids": [],
                "chunks_uploaded": 0,
                "images_uploaded": len(image_urls),
                "images_expected": images_expected,
                "image_urls": image_urls,
                "errors": errors,
                "toc_errors": [],
                "rearrange_notes": rearrange_notes,
                "rearrange_llm_raw": rearrange_llm_raw or "",
                "rearrange_llm_raw_chars": len(rearrange_llm_raw or ""),
                "embedding_used": use_emb,
                "vector_dim": dim,
            }
    else:
        _progress("images", 100, format_export_images_progress_label(done=0, total=0))

    chunks = chunk_text(md_body, chunk_size, chunk_overlap, split_text=split_text) or (
        [md_body] if md_body else []
    )
    total_chunks = len(chunks)
    if total_chunks == 0:
        _progress("embedding", 100, "no chunks to embed")
        _progress("db", 100, "no chunks to import")

    if chunks and use_emb:
        for idx, chunk in enumerate(chunks):
            content_hash = hashlib.md5(chunk.encode()).hexdigest()
            token_count = estimate_token_count(chunk)
            try:
                vec, embed_err = _embed_with_error(embeddings, chunk)
                if not vec:
                    errors.append(f"doc[{idx}]: embedding failed - {embed_err or 'unknown error'}")
                    continue
                _progress("embedding", int(((idx + 1) / total_chunks) * 100), f"embedded chunk {idx + 1}/{total_chunks}")
                vec_str = f"[{','.join(str(float(x)) for x in vec)}]"
                with conn.cursor() as cur:
                    cur.execute(
                        psql_sql.SQL(
                            """INSERT INTO {}
                           (content, document_name, chunk_index, token_count, created_by, updated_by, metadata, embedding)
                           VALUES (%s, %s, %s, %s, %s, %s, %s::jsonb, %s::vector) RETURNING id"""
                        ).format(psql_sql.Identifier(vec_tbl)),
                        (
                            chunk,
                            doc_name,
                            idx,
                            token_count,
                            actor_username,
                            actor_username,
                            json.dumps(
                                {
                                    "source_filename": src,
                                    "document_name": doc_name,
                                    "bucket_url": bucket_url,
                                    "source_file_url": source_file_url,
                                    "source_file_storage_path": source_storage_path,
                                    "chunk_index": idx,
                                    "chunk_type": "text",
                                    "token_count": token_count,
                                    "content_hash": content_hash,
                                },
                                ensure_ascii=False,
                            ),
                            vec_str,
                        ),
                    )
                    row_id = cur.fetchone()
                    if row_id:
                        doc_ids.append(row_id[0])
                conn.commit()
                _progress("db", int(((idx + 1) / total_chunks) * 100), f"imported chunk {idx + 1}/{total_chunks}")
            except Exception as e:
                conn.rollback()
                errors.append(f"doc[{idx}]: {e}")

        if uploaded_images:
            img_doc_ids = _insert_image_vector_rows(
                conn=conn,
                vec_tbl=vec_tbl,
                embeddings=embeddings,
                images=uploaded_images,
                base_chunk_index=total_chunks,
                doc_name=doc_name,
                src=src,
                bucket_url=bucket_url,
                source_file_url=source_file_url,
                source_storage_path=source_storage_path,
                actor_username=actor_username,
                errors=errors,
                progress_callback=_progress,
            )
            doc_ids.extend(img_doc_ids)

    elif chunks and not use_emb:
        _progress("embedding", 100, "embedding disabled")
        for idx, chunk in enumerate(chunks):
            content_hash = hashlib.md5(chunk.encode()).hexdigest()
            token_count = estimate_token_count(chunk)
            try:
                with conn.cursor() as cur:
                    cur.execute(
                        psql_sql.SQL(
                            """INSERT INTO {}
                           (content, document_name, chunk_index, token_count, created_by, updated_by, metadata)
                           VALUES (%s, %s, %s, %s, %s, %s, %s::jsonb) RETURNING id"""
                        ).format(psql_sql.Identifier(vec_tbl)),
                        (
                            chunk,
                            doc_name,
                            idx,
                            token_count,
                            actor_username,
                            actor_username,
                            json.dumps(
                                {
                                    "source_filename": src,
                                    "document_name": doc_name,
                                    "bucket_url": bucket_url,
                                    "source_file_url": source_file_url,
                                    "source_file_storage_path": source_storage_path,
                                    "chunk_index": idx,
                                    "token_count": token_count,
                                    "content_hash": content_hash,
                                },
                                ensure_ascii=False,
                            ),
                        ),
                    )
                    row_id = cur.fetchone()
                    if row_id:
                        doc_ids.append(row_id[0])
                conn.commit()
                _progress("db", int(((idx + 1) / total_chunks) * 100), f"imported chunk {idx + 1}/{total_chunks}")
            except Exception as e:
                conn.rollback()
                errors.append(f"doc[{idx}]: {e}")

    conn.close()

    toc_errors: list[str] = []

    imported_count = len(doc_ids)
    pending_summary = str(row.get("llm_summary") or "").strip()
    if imported_count > 0:
        from app.document_api.services.folder_catalog import ensure_folder_table, insert_folder_files

        ensure_folder_table()
        summary_map = (
            {(row["document_name"], src): pending_summary}
            if pending_summary
            else None
        )
        insert_folder_files(
            [(row["document_name"], src)],
            uploaded_by=row.get("created_by"),
            actor_username=actor_username,
            approved_by=actor_username,
            record_approval=bool(actor_username),
            llm_summaries=summary_map,
            department_id=row.get("department_id"),
        )
        delete_pending_row(pending_id)

    return {
        "status": "success" if not errors else ("partial" if imported_count > 0 else "error"),
        "pending_ingest_id": pending_id,
        "source_filename": src,
        "doc_ids": doc_ids,
        "chunks_uploaded": len(doc_ids),
        "images_uploaded": len(image_urls),
        "images_expected": images_expected,
        "image_urls": image_urls,
        "bucket_url": bucket_url,
        "source_file_url": source_file_url,
        "source_file_storage_path": source_storage_path,
        "errors": errors,
        "toc_errors": toc_errors,
        "embedding_used": use_emb,
        "vector_dim": dim,
        "table_name": sb.table_name,
        "query_name": sb.query_name,
    }


def _apply_reranker(
    rows: list[dict],
    *,
    query_text: str,
    reranker: Any | None,
    top_k: int,
) -> list[dict]:
    if not rows or reranker is None or not hasattr(reranker, "rerank"):
        return rows
    from app.document_api.lm_engine.qwen_vl_reranker import document_payload_for_rerank

    docs = [document_payload_for_rerank(r) for r in rows]
    try:
        scored = reranker.rerank(query=query_text, documents=docs, top_n=top_k)
    except Exception:
        return rows
    if not scored:
        return rows

    reordered: list[dict] = []
    for idx, rerank_score in scored:
        if 0 <= idx < len(rows):
            row = dict(rows[idx])
            row["rerank_score"] = rerank_score
            reordered.append(row)
    return reordered or rows


def query_documents(
    *,
    query_text: str,
    mode: str,
    pg: PgConfig,
    sb: SupabaseConfig,
    number_of_results: int | None = None,
    rrf_k: int = 60,
    filter_docs: list[str] | None = None,
    embeddings: Any | None = None,
    reranker: Any | None = None,
    rerank_candidates: int = 50,
) -> list[dict]:
    q = (query_text or "").strip()
    if not q:
        return []

    conn = _pg_conn(pg)
    effective_results = (
        number_of_results
        if number_of_results is not None
        else default_mcp_km_number_of_results()
    )
    top_k = max(1, effective_results)
    docs = [d.strip() for d in (filter_docs or []) if (d or "").strip()]
    docs_set = set(docs)
    rows: list[dict] = []

    def _row_source_filename(row: dict) -> str:
        if row.get("source_filename"):
            return str(row.get("source_filename") or "")
        metadata = row.get("metadata") or {}
        if isinstance(metadata, str):
            try:
                metadata = json.loads(metadata)
            except Exception:
                metadata = {}
        if isinstance(metadata, dict):
            return str(metadata.get("source_filename") or "")
        return ""

    try:
        from app.document_api.core.config import get_settings as _get_settings

        _q_settings = _get_settings()
        dim = resolve_vector_dim(
            embeddings,
            configured_dim=_q_settings.embedding_dimensions,
            use_embedding=embeddings is not None,
        )
        vec_tbl = _vector_table_ident(sb.table_name)
        setup_ddl(conn, dim, vec_tbl)
        ensure_schema(conn)
        filter_json = {}
        use_rerank = reranker is not None and mode in {"hybrid", "semantic"}
        recall_k = (
            max(top_k, rerank_candidates)
            if use_rerank
            else (top_k if not docs else max(top_k * 5, 50))
        )

        if mode == "hybrid":
            q_vec = _embed(embeddings, q) if embeddings else None
            if q_vec is None:
                mode = "keyword"
            else:
                fetch_k = recall_k if use_rerank else (top_k if not docs else max(top_k * 5, 50))
                vec_str = f"[{','.join(str(round(x, 8)) for x in q_vec)}]"
                with conn.cursor() as cur:
                    cur.execute(
                        """SELECT id, content, metadata, similarity,
                                  similarity AS hybrid_score,
                                  similarity AS semantic_score,
                                  0.0::FLOAT AS keyword_rank
                           FROM match_documents(%s::vector, %s, %s::jsonb)""",
                        (vec_str, fetch_k, json.dumps(filter_json, ensure_ascii=False)),
                    )
                    cols = [d[0] for d in cur.description]
                    rows = [dict(zip(cols, r)) for r in cur.fetchall()]
                if docs_set:
                    rows = [r for r in rows if _row_source_filename(r) in docs_set]
                if use_rerank:
                    rows = _apply_reranker(rows, query_text=q, reranker=reranker, top_k=top_k)
                else:
                    rows = rows[:top_k]

        if mode == "semantic":
            if embeddings is None:
                mode = "keyword"
            else:
                try:
                    from langchain_community.vectorstores import SupabaseVectorStore

                    sb_client = _sb_client(sb)
                    vs = SupabaseVectorStore(
                        client=sb_client,
                        embedding=embeddings,
                        table_name=sb.table_name,
                        query_name=sb.query_name,
                    )
                    fetch_k = recall_k if use_rerank else (top_k if not docs else max(top_k * 5, 50))
                    search_kwargs: dict = {"k": fetch_k}
                    scored = vs.similarity_search_with_relevance_scores(query=q, **search_kwargs)
                    rows = [
                        {
                            "id": None,
                            "source_filename": doc.metadata.get("source_filename", ""),
                            "chunk_index": doc.metadata.get("chunk_index", 0),
                            "content": doc.page_content,
                            "similarity": float(score),
                        }
                        for doc, score in scored
                    ]
                    if docs_set:
                        rows = [r for r in rows if r.get("source_filename") in docs_set]
                    if use_rerank:
                        rows = _apply_reranker(rows, query_text=q, reranker=reranker, top_k=top_k)
                    else:
                        rows = rows[:top_k]
                except Exception:
                    q_vec = _embed(embeddings, q)
                    if q_vec:
                        fetch_k = recall_k if use_rerank else (top_k if not docs else max(top_k * 5, 50))
                        vec_str = f"[{','.join(str(round(x, 8)) for x in q_vec)}]"
                        with conn.cursor() as cur:
                            cur.execute(
                                """SELECT id, content, metadata, similarity
                                   FROM match_documents(%s::vector, %s, %s::jsonb)""",
                                (vec_str, fetch_k, json.dumps(filter_json, ensure_ascii=False)),
                            )
                            cols = [d[0] for d in cur.description]
                            rows = [dict(zip(cols, r)) for r in cur.fetchall()]
                        if docs_set:
                            rows = [r for r in rows if _row_source_filename(r) in docs_set]
                        if use_rerank:
                            rows = _apply_reranker(rows, query_text=q, reranker=reranker, top_k=top_k)
                        else:
                            rows = rows[:top_k]

                if not rows:
                    keywords = [w for w in q.split() if len(w) >= 2]
                    if keywords:
                        like_clause = " OR ".join(["content ILIKE %s"] * len(keywords))
                        filter_clause = ""
                        params: list[Any] = [f"%{k}%" for k in keywords]
                        if docs:
                            filter_clause = " AND metadata->>'source_filename' = ANY(%s)"
                            params.append(docs)
                        with conn.cursor() as cur:
                            cur.execute(
                                f"""SELECT id, content, metadata,
                                           0.0::FLOAT AS similarity
                                    FROM {vec_tbl}
                                    WHERE ({like_clause}){filter_clause}
                                    LIMIT %s""",
                                params + [top_k],
                            )
                            cols = [d[0] for d in cur.description]
                            rows = [dict(zip(cols, r)) for r in cur.fetchall()]

        if mode == "keyword":
            with conn.cursor() as cur:
                cur.execute(
                    f"""SELECT id, content, metadata, 0.0::FLOAT AS rank
                       FROM {vec_tbl}
                       WHERE content ILIKE %s
                         AND (%s::text[] IS NULL OR metadata->>'source_filename' = ANY(%s))
                       LIMIT %s""",
                    (f"%{q}%", docs or None, docs or None, top_k),
                )
                cols = [d[0] for d in cur.description]
                rows = [dict(zip(cols, r)) for r in cur.fetchall()]

            if not rows:
                keywords = [w for w in q.split() if len(w) >= 2]
                if keywords:
                    like_clause = " OR ".join(["content ILIKE %s"] * len(keywords))
                    filter_clause = ""
                    params = [f"%{k}%" for k in keywords]
                    if docs:
                        filter_clause = " AND metadata->>'source_filename' = ANY(%s)"
                        params.append(docs)
                    with conn.cursor() as cur:
                        cur.execute(
                            f"""SELECT id, content, metadata,
                                       0.0::FLOAT AS rank
                                FROM {vec_tbl}
                                WHERE ({like_clause}){filter_clause}
                                LIMIT %s""",
                            params + [top_k],
                        )
                        cols = [d[0] for d in cur.description]
                        rows = [dict(zip(cols, r)) for r in cur.fetchall()]

    finally:
        conn.close()

    from app.document_api.core.storage_urls import rewrite_storage_urls_in_row

    for i, row in enumerate(rows):
        metadata = row.get("metadata") or {}
        if isinstance(metadata, str):
            try:
                metadata = json.loads(metadata)
            except Exception:
                metadata = {}
        if "source_filename" not in row:
            row["source_filename"] = metadata.get("source_filename", "")
        if "chunk_index" not in row:
            row["chunk_index"] = metadata.get("chunk_index", 0)
        for k, v in list(row.items()):
            if hasattr(v, "__float__"):
                row[k] = float(v)
        rows[i] = rewrite_storage_urls_in_row(row)
    return rows

